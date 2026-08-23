"""
Extract text from study documents.

Supports:
- PDF
- Scanned PDF
- Handwritten PDF
- DOCX
- PPTX
- TXT
- JPG / JPEG
- PNG
- WEBP
- GIF
- BMP

Gemini is used for OCR of:
- scanned PDFs
- handwritten PDFs
- images
- screenshots
- handwritten notes
"""

from io import BytesIO
from pathlib import Path
import os

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from google import genai
from google.genai import types

# PyMuPDF is used to render scanned PDF pages as images
import fitz


MAX_CHARS = 60_000

OCR_MODEL = os.getenv(
    "GEMINI_MODEL",
    "gemini-2.0-flash",
)


class UnsupportedDocumentError(Exception):
    pass


# ============================================================
# GEMINI CLIENT
# ============================================================

def _get_gemini_client():
    api_key = os.getenv("GEMINI_API_KEY")

    if not api_key:
        raise RuntimeError(
            "GEMINI_API_KEY is not configured."
        )

    return genai.Client(
        api_key=api_key
    )


# ============================================================
# MAIN EXTRACTION FUNCTION
# ============================================================

def extract_text(
    data: bytes,
    mime_type: str,
    filename: str = "document",
) -> str:

    suffix = Path(filename).suffix.lower()

    mime_type = (
        mime_type or ""
    ).lower().split(";")[0].strip()

    # ---------------------------------------------------------
    # PDF
    # ---------------------------------------------------------

    if (
        mime_type == "application/pdf"
        or suffix == ".pdf"
    ):
        text = _extract_pdf(data)

        # If normal PDF extraction fails,
        # render pages and use Gemini OCR.
        if not text.strip():
            text = _ocr_pdf(data)

    # ---------------------------------------------------------
    # DOCX
    # ---------------------------------------------------------

    elif (
        mime_type
        == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        or suffix == ".docx"
    ):
        text = _extract_docx(data)

    # ---------------------------------------------------------
    # PPTX
    # ---------------------------------------------------------

    elif (
        mime_type
        == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        or suffix == ".pptx"
    ):
        text = _extract_pptx(data)

    # ---------------------------------------------------------
    # TXT
    # ---------------------------------------------------------

    elif (
        mime_type == "text/plain"
        or suffix == ".txt"
    ):
        text = _extract_txt(data)

    # ---------------------------------------------------------
    # IMAGES
    # ---------------------------------------------------------

    elif (
        mime_type.startswith("image/")
        or suffix in {
            ".jpg",
            ".jpeg",
            ".png",
            ".webp",
            ".gif",
            ".bmp",
        }
    ):
        text = _ocr_image(
            data,
            mime_type,
            suffix,
        )

    else:
        raise UnsupportedDocumentError(
            f"Unsupported document type: "
            f"{mime_type or suffix}"
        )

    text = text.strip()

    if not text:
        raise ValueError(
            "No readable text was found in this document. "
            "Try uploading a clearer PDF or image."
        )

    if len(text) > MAX_CHARS:
        text = (
            text[:MAX_CHARS]
            + "\n\n[Document truncated for length]"
        )

    return text


# ============================================================
# NORMAL PDF TEXT EXTRACTION
# ============================================================

def _extract_pdf(data: bytes) -> str:

    reader = PdfReader(
        BytesIO(data)
    )

    parts = []

    for page in reader.pages:

        page_text = (
            page.extract_text()
            or ""
        )

        if page_text.strip():
            parts.append(page_text)

    return "\n\n".join(parts)


# ============================================================
# OCR SCANNED / HANDWRITTEN PDF
# ============================================================

def _ocr_pdf(data: bytes) -> str:
    """
    Render each PDF page as an image and send it to Gemini.

    Supports:
    - scanned PDFs
    - handwritten PDFs
    - photos saved as PDFs
    - screenshots saved as PDFs
    """

    pdf = fitz.open(
        stream=data,
        filetype="pdf",
    )

    parts = []

    try:

        for page_number, page in enumerate(
            pdf
        ):

            # Render page at good resolution
            pix = page.get_pixmap(
                matrix=fitz.Matrix(2, 2),
                alpha=False,
            )

            image_bytes = pix.tobytes(
                "png"
            )

            page_text = _ocr_image_bytes(
                image_bytes,
                "image/png",
                page_number + 1,
            )

            if page_text.strip():

                parts.append(
                    f"[Page {page_number + 1}]\n"
                    f"{page_text}"
                )

            # Avoid processing extremely
            # large documents
            if (
                sum(len(x) for x in parts)
                >= MAX_CHARS
            ):
                break

    finally:
        pdf.close()

    return "\n\n".join(parts)


# ============================================================
# IMAGE OCR
# ============================================================

def _ocr_image(
    data: bytes,
    mime_type: str,
    suffix: str = "",
) -> str:

    if not mime_type.startswith(
        "image/"
    ):

        mime_type = {
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".webp": "image/webp",
            ".gif": "image/gif",
            ".bmp": "image/bmp",
        }.get(
            suffix,
            "image/png",
        )

    return _ocr_image_bytes(
        data,
        mime_type,
        1,
    )


# ============================================================
# GEMINI IMAGE OCR
# ============================================================

def _ocr_image_bytes(
    image_bytes: bytes,
    mime_type: str,
    page_number: int = 1,
) -> str:

    client = _get_gemini_client()

    prompt = """
You are a highly accurate OCR system for QuizNest.

Extract ALL readable educational content from
the provided image.

The image may contain:

- printed text
- handwritten text
- handwritten notes
- questions
- multiple-choice options
- headings
- paragraphs
- mathematical expressions
- tables
- diagrams containing readable text
- screenshots
- classroom notes

IMPORTANT RULES:

1. Transcribe the content accurately.
2. Preserve the original meaning.
3. Preserve headings and question structure.
4. Include handwritten and printed text.
5. Do not summarize.
6. Do not explain.
7. Do not invent missing content.
8. If something is genuinely unreadable, write [unclear].
9. Keep mathematical expressions as accurately as possible.
10. Return ONLY the extracted/transcribed text.
"""

    image_part = types.Part.from_bytes(
        data=image_bytes,
        mime_type=mime_type,
    )

    response = client.models.generate_content(
        model=OCR_MODEL,
        contents=[
            image_part,
            prompt,
        ],
        config=types.GenerateContentConfig(
            temperature=0,
        ),
    )

    if not response.text:
        return ""

    return response.text.strip()


# ============================================================
# DOCX
# ============================================================

def _extract_docx(
    data: bytes,
) -> str:

    doc = DocxDocument(
        BytesIO(data)
    )

    parts = []

    # Normal paragraphs
    for paragraph in doc.paragraphs:

        text = paragraph.text.strip()

        if text:
            parts.append(text)

    # Tables
    for table in doc.tables:

        for row in table.rows:

            row_text = " | ".join(
                cell.text.strip()
                for cell in row.cells
            )

            if row_text.strip():
                parts.append(row_text)

    return "\n".join(parts)


# ============================================================
# TXT
# ============================================================

def _extract_txt(
    data: bytes,
) -> str:

    return data.decode(
        "utf-8",
        errors="ignore",
    )


# ============================================================
# PPTX
# ============================================================

def _extract_pptx(
    data: bytes,
) -> str:

    prs = Presentation(
        BytesIO(data)
    )

    parts = []

    for slide_number, slide in enumerate(
        prs.slides,
        start=1,
    ):

        slide_parts = []

        for shape in slide.shapes:

            if shape.has_text_frame:

                text = shape.text.strip()

                if text:
                    slide_parts.append(
                        text
                    )

        if slide_parts:

            parts.append(
                f"[Slide {slide_number}]\n"
                + "\n".join(slide_parts)
            )

    return "\n\n".join(parts)